import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import asyncio
from dataclasses import dataclass, asdict
from typing import Dict, List, Any, Optional
from datetime import datetime
import uuid
import json
import tempfile
from pptx import Presentation
import pandas as pd
from docx import Document

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# NEW: Model Context Protocol (MCP) Implementation
@dataclass
class MCPMessage:
    """Model Context Protocol message structure"""
    type: str
    sender: str
    receiver: str
    trace_id: str
    timestamp: str
    payload: Dict[str, Any]
    
    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)

class MCPMessageBus:
    """Message bus for agent communication"""
    def __init__(self):
        self.messages: List[MCPMessage] = []
        self.subscribers: Dict[str, List[callable]] = {}
    
    def subscribe(self, agent_name: str, callback: callable):
        if agent_name not in self.subscribers:
            self.subscribers[agent_name] = []
        self.subscribers[agent_name].append(callback)
    
    async def send_message(self, message: MCPMessage):
        self.messages.append(message)
        if message.receiver in self.subscribers:
            for callback in self.subscribers[message.receiver]:
                await callback(message)

# NEW: Base Agent Class
class BaseAgent:
    """Base class for all agents"""
    def __init__(self, name: str, message_bus: MCPMessageBus):
        self.name = name
        self.message_bus = message_bus
        self.message_bus.subscribe(name, self.handle_message)
    
    async def handle_message(self, message: MCPMessage):
        pass
    
    async def send_message(self, receiver: str, message_type: str, payload: Dict[str, Any], trace_id: str):
        message = MCPMessage(
            type=message_type,
            sender=self.name,
            receiver=receiver,
            trace_id=trace_id,
            timestamp=datetime.now().isoformat(),
            payload=payload
        )
        await self.message_bus.send_message(message)

# NEW: Agent 1 - Ingestion Agent (Enhanced from your existing PDF logic)
class IngestionAgent(BaseAgent):
    """Handles document parsing and preprocessing - Enhanced from your get_pdf_text function"""
    
    def __init__(self, message_bus: MCPMessageBus):
        super().__init__("IngestionAgent", message_bus)
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=500)
    
    async def handle_message(self, message: MCPMessage):
        if message.type == "INGEST_DOCUMENT":
            await self.process_document(message)
    
    async def process_document(self, message: MCPMessage):
        file_path = message.payload.get("file_path")
        file_name = message.payload.get("file_name")
        trace_id = message.trace_id
        
        try:
            # Determine file type and parse accordingly
            _, ext = os.path.splitext(file_name.lower())
            
            if ext == '.pdf':
                content = self.get_pdf_text([file_path])  # Using your existing function logic
            elif ext == '.pptx':
                content = self.parse_pptx(file_path)
            elif ext == '.csv':
                content = self.parse_csv(file_path)
            elif ext == '.docx':
                content = self.parse_docx(file_path)
            elif ext in ['.txt', '.md']:
                content = self.parse_text(file_path)
            else:
                raise ValueError(f"Unsupported file format: {ext}")
            
            # Use your existing text chunking logic
            chunks = self.get_text_chunks(content)
            
            # Send to RetrievalAgent
            await self.send_message(
                receiver="RetrievalAgent",
                message_type="CONTENT_PROCESSED",
                payload={
                    "content": content,
                    "file_name": file_name,
                    "chunks": chunks
                },
                trace_id=trace_id
            )
            
        except Exception as e:
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="ERROR",
                payload={"error": str(e), "stage": "ingestion"},
                trace_id=trace_id
            )
    
    def get_pdf_text(self, pdf_docs):
        """Your existing PDF processing function"""
        text = ""
        for pdf in pdf_docs:
            if isinstance(pdf, str):  # File path
                with open(pdf, 'rb') as file:
                    pdf_reader = PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text()
            else:  # File object
                pdf_reader = PdfReader(pdf)
                for page in pdf_reader.pages:
                    text += page.extract_text()
        return text
    
    def get_text_chunks(self, text):
        """Your existing text chunking function"""
        chunks = self.text_splitter.split_text(text)
        return chunks
    
    # NEW: Additional document format parsers
    def parse_pptx(self, file_path: str) -> str:
        """NEW: PowerPoint parsing"""
        presentation = Presentation(file_path)
        content = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    content += shape.text + "\n"
        return content
    
    def parse_csv(self, file_path: str) -> str:
        """NEW: CSV parsing"""
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def parse_docx(self, file_path: str) -> str:
        """NEW: Word document parsing"""
        doc = Document(file_path)
        content = ""
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
        return content
    
    def parse_text(self, file_path: str) -> str:
        """NEW: Text/Markdown parsing"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

# NEW: Agent 2 - Retrieval Agent (Enhanced from your existing FAISS logic)
class RetrievalAgent(BaseAgent):
    """Handles embeddings and retrieval - Enhanced from your get_vector_store function"""
    
    def __init__(self, message_bus: MCPMessageBus):
        super().__init__("RetrievalAgent", message_bus)
        self.embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        self.vector_stores = {}  # Store multiple vector stores for different documents
    
    async def handle_message(self, message: MCPMessage):
        if message.type == "CONTENT_PROCESSED":
            await self.store_embeddings(message)
        elif message.type == "RETRIEVE_CONTEXT":
            await self.retrieve_context(message)
    
    async def store_embeddings(self, message: MCPMessage):
        chunks = message.payload.get("chunks", [])
        file_name = message.payload.get("file_name")
        trace_id = message.trace_id
        
        try:
            # Use your existing vector store logic
            vector_store = self.get_vector_store(chunks)
            self.vector_stores[file_name] = vector_store
            
            # Store in session state for UI compatibility
            if "vector_store" not in st.session_state:
                st.session_state["vector_store"] = vector_store
            else:
                # Merge with existing vector store if needed
                st.session_state["vector_store"] = vector_store
            
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="EMBEDDINGS_STORED",
                payload={"file_name": file_name, "chunks_count": len(chunks)},
                trace_id=trace_id
            )
            
        except Exception as e:
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="ERROR",
                payload={"error": str(e), "stage": "retrieval_storage"},
                trace_id=trace_id
            )
    
    def get_vector_store(self, text_chunks):
        """Your existing vector store function"""
        vector_store = FAISS.from_texts(text_chunks, embedding=self.embeddings)
        return vector_store
    
    async def retrieve_context(self, message: MCPMessage):
        query = message.payload.get("query")
        top_k = message.payload.get("top_k", 4)
        trace_id = message.trace_id
        
        try:
            # Use your existing similarity search logic
            if "vector_store" in st.session_state:
                docs = st.session_state["vector_store"].similarity_search(query, k=top_k)
                
                retrieved_chunks = []
                for doc in docs:
                    retrieved_chunks.append({
                        "content": doc.page_content,
                        "metadata": doc.metadata if hasattr(doc, 'metadata') else {}
                    })
                
                await self.send_message(
                    receiver="LLMResponseAgent",
                    message_type="CONTEXT_RETRIEVED",
                    payload={
                        "retrieved_context": retrieved_chunks,
                        "query": query,
                        "docs": docs  # Pass original docs for chain compatibility
                    },
                    trace_id=trace_id
                )
            else:
                raise ValueError("No vector store available")
                
        except Exception as e:
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="ERROR",
                payload={"error": str(e), "stage": "retrieval_query"},
                trace_id=trace_id
            )

# NEW: Agent 3 - LLM Response Agent (Enhanced from your existing chain logic)
class LLMResponseAgent(BaseAgent):
    """Handles LLM response generation - Enhanced from your get_conversational_chain function"""
    
    def __init__(self, message_bus: MCPMessageBus):
        super().__init__("LLMResponseAgent", message_bus)
        self.chain = self.get_conversational_chain()
    
    async def handle_message(self, message: MCPMessage):
        if message.type == "CONTEXT_RETRIEVED":
            await self.generate_response(message)
    
    def get_conversational_chain(self):
        """Your existing conversational chain function"""
        prompt_template = """
        Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
        provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
        Context:\n {context}?\n
        Question: \n{question}\n

        Answer:
        """
        
        model = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.3)
        prompt = PromptTemplate(template=prompt_template, input_variables=['context', 'question'])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    
    async def generate_response(self, message: MCPMessage):
        retrieved_context = message.payload.get("retrieved_context", [])
        query = message.payload.get("query")
        docs = message.payload.get("docs", [])
        trace_id = message.trace_id
        
        try:
            # Use your existing chain logic
            response = self.chain(
                {"input_documents": docs, "question": query},
                return_only_outputs=True
            )
            
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="RESPONSE_GENERATED",
                payload={
                    "response": response["output_text"],
                    "query": query,
                    "sources": [chunk['metadata'] for chunk in retrieved_context]
                },
                trace_id=trace_id
            )
            
        except Exception as e:
            await self.send_message(
                receiver="CoordinatorAgent",
                message_type="ERROR",
                payload={"error": str(e), "stage": "llm_response"},
                trace_id=trace_id
            )

# NEW: Agent 4 - Coordinator Agent (NEW - orchestrates everything)
class CoordinatorAgent(BaseAgent):
    """Orchestrates the entire workflow"""
    
    def __init__(self, message_bus: MCPMessageBus):
        super().__init__("CoordinatorAgent", message_bus)
        self.active_traces: Dict[str, Dict[str, Any]] = {}
    
    async def handle_message(self, message: MCPMessage):
        if message.type in ["EMBEDDINGS_STORED", "RESPONSE_GENERATED", "ERROR"]:
            await self.handle_workflow_completion(message)
    
    async def process_document_upload(self, file_path: str, file_name: str) -> str:
        """Process uploaded document"""
        trace_id = str(uuid.uuid4())
        
        self.active_traces[trace_id] = {
            "status": "processing",
            "file_name": file_name,
            "start_time": datetime.now()
        }
        
        await self.send_message(
            receiver="IngestionAgent",
            message_type="INGEST_DOCUMENT",
            payload={"file_path": file_path, "file_name": file_name},
            trace_id=trace_id
        )
        
        return trace_id
    
    async def process_query(self, query: str) -> str:
        """Process user query"""
        trace_id = str(uuid.uuid4())
        
        self.active_traces[trace_id] = {
            "status": "querying",
            "query": query,
            "start_time": datetime.now()
        }
        
        await self.send_message(
            receiver="RetrievalAgent",
            message_type="RETRIEVE_CONTEXT",
            payload={"query": query, "top_k": 4},
            trace_id=trace_id
        )
        
        return trace_id
    
    async def handle_workflow_completion(self, message: MCPMessage):
        trace_id = message.trace_id
        if trace_id in self.active_traces:
            if message.type == "RESPONSE_GENERATED":
                self.active_traces[trace_id]["status"] = "completed"
                self.active_traces[trace_id]["response"] = message.payload
            elif message.type == "EMBEDDINGS_STORED":
                self.active_traces[trace_id]["status"] = "stored"
            elif message.type == "ERROR":
                self.active_traces[trace_id]["status"] = "error"
                self.active_traces[trace_id]["error"] = message.payload
    
    def get_trace_status(self, trace_id: str) -> Optional[Dict[str, Any]]:
        return self.active_traces.get(trace_id)

# NEW: Initialize the agentic system
@st.cache_resource
def initialize_agentic_system():
    """Initialize the multi-agent system"""
    message_bus = MCPMessageBus()
    
    agents = {
        "coordinator": CoordinatorAgent(message_bus),
        "ingestion": IngestionAgent(message_bus),
        "retrieval": RetrievalAgent(message_bus),
        "llm_response": LLMResponseAgent(message_bus)
    }
    
    return agents, message_bus

# ENHANCED: Main function with agentic capabilities
def main():
    st.set_page_config("Agentic Chat with Multiple Documents 🤖", layout="wide")
    st.header("Agentic RAG Chatbot with MCP - Multi-Format Document QA 💁")
    
    # Initialize agentic system
    agents, message_bus = initialize_agentic_system()
    coordinator = agents["coordinator"]
    
    # Enhanced sidebar with multi-format support
    with st.sidebar:
        st.title("📁 Document Upload Menu")
        
        # Multi-format file uploader
        uploaded_files = st.file_uploader(
            "Upload your Documents (PDF, PPTX, CSV, DOCX, TXT, MD)",
            accept_multiple_files=True,
            type=['pdf', 'pptx', 'csv', 'docx', 'txt', 'md']
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.write(f"📄 {uploaded_file.name}")
                with col2:
                    if st.button("Process", key=f"process_{uploaded_file.name}"):
                        with st.spinner(f"Processing {uploaded_file.name}..."):
                            # Save uploaded file temporarily
                            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
                                tmp_file.write(uploaded_file.getvalue())
                                tmp_file_path = tmp_file.name
                            
                            # Process document using agentic system
                            loop = asyncio.new_event_loop()
                            asyncio.set_event_loop(loop)
                            trace_id = loop.run_until_complete(
                                coordinator.process_document_upload(tmp_file_path, uploaded_file.name)
                            )
                            
                            # Wait for processing to complete
                            import time
                            for _ in range(30):  # Wait up to 30 seconds
                                time.sleep(1)
                                status = coordinator.get_trace_status(trace_id)
                                if status and status.get("status") == "stored":
                                    st.success(f"✅ {uploaded_file.name} processed!")
                                    break
                                elif status and status.get("status") == "error":
                                    st.error(f"❌ Error: {status.get('error', {}).get('error', 'Unknown error')}")
                                    break
                            
                            # Clean up
                            os.unlink(tmp_file_path)
        
        # NEW: MCP Message Monitor
        st.subheader("🔧 System Monitor")
        if st.button("Show MCP Messages"):
            with st.expander("Recent MCP Messages"):
                recent_messages = message_bus.messages[-10:] if message_bus.messages else []
                for msg in recent_messages:
                    st.json({
                        "type": msg.type,
                        "sender": msg.sender,
                        "receiver": msg.receiver,
                        "timestamp": msg.timestamp,
                        "trace_id": msg.trace_id,
                        "payload": msg.payload 
                    })
    
    # Enhanced chat interface
    st.subheader("💬 Chat Interface")
    
    # Initialize chat history
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    
    # Display chat history
    for message in st.session_state.chat_history:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            # if "sources" in message and message["sources"]:
            #     with st.expander("📚 Sources"):
            #         for source in message["sources"]:
            #             st.write(f"- {source}")
    
    # Chat input
    user_question = st.chat_input("Ask a question about your documents...")
    
    if user_question:
        # Add user message to chat history
        st.session_state.chat_history.append({"role": "user", "content": user_question})
        
        # Display user message
        with st.chat_message("user"):
            st.markdown(user_question)
        
        # Process query using agentic system
        with st.chat_message("assistant"):
            with st.spinner("🤖 Agents working..."):
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                trace_id = loop.run_until_complete(coordinator.process_query(user_question))
                
                # Wait for response
                import time
                response_content = None
                sources = []
                
                for _ in range(60):  # Wait up to 60 seconds
                    time.sleep(1)
                    status = coordinator.get_trace_status(trace_id)
                    if status and status.get("status") == "completed":
                        response_data = status.get("response", {})
                        response_content = response_data.get("response", "No response generated")
                        sources = response_data.get("sources", [])
                        break
                    elif status and status.get("status") == "error":
                        response_content = f"❌ Error: {status.get('error', {}).get('error', 'Unknown error')}"
                        break
                
                if response_content is None:
                    response_content = "⚠️ Sorry, I couldn't generate a response. Please try again."
                
                st.markdown(response_content)
                
                if sources:
                    with st.expander("📚 Sources"):
                        for source in sources:
                            st.write(f"- {source}")
        
        # Add assistant response to chat history
        st.session_state.chat_history.append({
            "role": "assistant",
            "content": response_content,
            "sources": sources
        })
    
    # Enhanced system information
    with st.expander("🔍 Agentic System Information"):
        col1, col2 = st.columns(2)
        
        with col1:
            st.write("**Active Agents:**")
            st.write("- 🎯 CoordinatorAgent: Orchestrates workflow")
            st.write("- 📄 IngestionAgent: Processes documents")
            st.write("- 🔍 RetrievalAgent: Handles embeddings")
            st.write("- 🤖 LLMResponseAgent: Generates responses")
        
        with col2:
            st.write("**Supported Formats:**")
            st.write("- 📄 PDF Documents")
            st.write("- 📊 PowerPoint (PPTX)")
            st.write("- 📈 CSV Files")
            st.write("- 📝 Word Documents (DOCX)")
            st.write("- 📋 Text/Markdown Files")
        
        st.write(f"**MCP Messages:** {len(message_bus.messages)}")
        st.write(f"**Vector Store Status:** {'✅ Ready' if 'vector_store' in st.session_state else '❌ Not Ready'}")

if __name__ == "__main__":
    main()